import argparse
from pathlib import Path
from typing import Iterable, List

try:
    from pptx import Presentation  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
except ImportError as e:
    # Provide a clearer error when dependency is missing.
    raise SystemExit(
        "python-pptx が見つかりません。\n"
        "インストール例: pip install python-pptx または uv add python-pptx"
    ) from e


def _iter_shape_text(shape) -> Iterable[str]:
    """Yield text content from a shape, descending into groups and tables.

    - Handles text frames, tables, and grouped shapes recursively.
    - Strips surrounding whitespace and skips empty lines.
    """
    # Grouped shapes
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            yield from _iter_shape_text(sub)
        return

    # Tables
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                text = (cell.text or "").strip()
                if text:
                    for line in text.splitlines():
                        line = line.strip()
                        if line:
                            yield line
        return

    # Text frames
    if getattr(shape, "has_text_frame", False):
        text = (shape.text or "").strip()
        if text:
            for line in text.splitlines():
                line = line.strip()
                if line:
                    yield line


def extract_slide_lines(slide) -> List[str]:
    """Collect unique non-empty text lines from a slide.

    Deduplicates while preserving order (first occurrence wins).
    """
    seen = set()
    lines: List[str] = []

    # Prefer the slide title first (if present)
    title_text = None
    try:
        if getattr(slide.shapes, "title", None) is not None:
            title_text = (slide.shapes.title.text or "").strip()
    except Exception:
        title_text = None
    if title_text:
        for line in title_text.splitlines():
            line = line.strip()
            if line and line not in seen:
                seen.add(line)
                lines.append(line)

    # Other shapes
    for shape in slide.shapes:
        try:
            for line in _iter_shape_text(shape):
                if line not in seen:
                    seen.add(line)
                    lines.append(line)
        except Exception:
            # Skip shapes we cannot parse
            continue

    return lines


def extract_notes_lines(slide) -> List[str]:
    """Extract speaker notes lines if available."""
    try:
        notes_slide = getattr(slide, "notes_slide", None)
        if notes_slide and notes_slide.notes_text_frame:
            text = (notes_slide.notes_text_frame.text or "").strip()
            if text:
                return [ln.strip() for ln in text.splitlines() if ln.strip()]
    except Exception:
        pass
    return []


def to_markdown(pptx_path: Path, out_path: Path, split_files: bool = False) -> None:
    prs = Presentation(str(pptx_path))

    if split_files:
        out_dir = out_path
        out_dir.mkdir(parents=True, exist_ok=True)
    else:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        md = []
        md.append(f"# {pptx_path.name} 文字抽出\n")

    for idx, slide in enumerate(prs.slides, start=1):
        lines = extract_slide_lines(slide)
        notes = extract_notes_lines(slide)

        # Derive a slide heading: prefer first line as title fallback
        heading = f"スライド {idx}"
        if lines:
            heading = f"スライド {idx}: {lines[0]}"

        if split_files:
            slide_md = [f"# {heading}\n"]
            # skip duplicate of title line in bullets if we used it
            content_lines = lines[1:] if lines else []
            for line in content_lines:
                slide_md.append(f"- {line}")
            if notes:
                slide_md.append("\n## ノート")
                for n in notes:
                    slide_md.append(f"- {n}")
            slide_md.append("")

            slide_file = out_dir / f"slide_{idx:02d}.md"
            slide_file.write_text("\n".join(slide_md), encoding="utf-8")
        else:
            md.append(f"\n## {heading}\n")
            content_lines = lines[1:] if lines else []
            for line in content_lines:
                md.append(f"- {line}")
            if notes:
                md.append("\n### ノート")
                for n in notes:
                    md.append(f"- {n}")

    if not split_files:
        md.append("")  # trailing newline
        out_path.write_text("\n".join(md), encoding="utf-8")


def main():
    default_input = Path("../example/assets/自作Devinで全員がエンジニアに (1).pptx")
    default_output = Path("example/自作Devinで全員がエンジニアに (1).md")

    parser = argparse.ArgumentParser(
        description="PPTX から各スライドの文字を抽出して Markdown にまとめます。",
    )
    parser.add_argument(
        "--input",
        "-i",
        type=Path,
        default=default_input,
        help=f"入力 PPTX ファイルパス (既定: {default_input})",
    )
    parser.add_argument(
        "--output",
        "-o",
        type=Path,
        default=default_output,
        help=(
            "出力先: 単一 Markdown ファイルのパス。"
            f" --split 指定時は出力ディレクトリ (既定: {default_output})"
        ),
    )
    parser.add_argument(
        "--split",
        action="store_true",
        help="各スライドを個別の Markdown ファイルに分割して出力",
    )

    args = parser.parse_args()

    if not args.input.exists():
        raise SystemExit(f"入力ファイルが見つかりません: {args.input}")
    out_path = args.output
    # --split 指定時に .md を付けたままでも、拡張子を除いたディレクトリに出力する
    if args.split and out_path.suffix.lower() == ".md":
        out_path = out_path.with_suffix("")

    to_markdown(args.input, out_path, split_files=args.split)


if __name__ == "__main__":
    main()
